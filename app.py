import base64
import os
import requests
from dotenv import load_dotenv
import faiss
import numpy as np
from PIL import Image
from pypdf import PdfReader
from docx import Document  # python-docx
from pptx import Presentation  # python-pptx
from sentence_transformers import SentenceTransformer
import streamlit as st
from streamlit_pdf_viewer import pdf_viewer

# Prevent CPU thrashing by limiting backend parallel thread over-allocation
os.environ["OMP_NUM_THREADS"] = "4"
os.environ["MKL_NUM_THREADS"] = "4"

load_dotenv()
API_KEY = os.getenv("OPENROUTER_API_KEY")

st.set_page_config(
    page_title="Universal Multimodal AI Workspace",
    page_icon="🤖",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom injection styling
st.markdown("""
    <style>
        .block-container { padding-top: 2rem; }
        .stButton>button { width: 100%; border-radius: 6px; }
        .asset-card {
            background-color: #f8f9fa;
            border-left: 5px solid #ff4b4b;
            padding: 15px;
            border-radius: 4px;
            margin-bottom: 15px;
        }
    </style>
""", unsafe_allow_html=True)

@st.cache_resource
def load_embedding_model():
    return SentenceTransformer("all-MiniLM-L6-v2", model_kwargs={"torch_dtype": "float16"})

embed_model = load_embedding_model()

# Initialize Global Session States
if "text_chunks" not in st.session_state: st.session_state.text_chunks = []
if "faiss_index" not in st.session_state: st.session_state.faiss_index = None
if "doc_raw_text" not in st.session_state: st.session_state.doc_raw_text = ""
if "file_raw_bytes" not in st.session_state: st.session_state.file_raw_bytes = None
if "active_image_b64" not in st.session_state: st.session_state.active_image_b64 = None
if "current_file_ext" not in st.session_state: st.session_state.current_file_ext = ""
if "current_file_name" not in st.session_state: st.session_state.current_file_name = ""

# ---- CHAT MEMORY STRUCTURAL STORAGE ----
if "chat_history" not in st.session_state: st.session_state.chat_history = []
if "image_chat_history" not in st.session_state: st.session_state.image_chat_history = []
if "saved_doc_sessions" not in st.session_state: st.session_state.saved_doc_sessions = {}
if "saved_img_sessions" not in st.session_state: st.session_state.saved_img_sessions = {}

# --- CORE UTILITY FUNCTIONS ---
def chunk_text(text, chunk_size=500):
    chunks = []
    words = text.split()
    for i in range(0, len(words), chunk_size):
        chunks.append(" ".join(words[i:i + chunk_size]))
    return chunks

def build_vector_db(chunks):
    if not chunks: return None
    embeddings = embed_model.encode(chunks)
    embeddings_array = np.array(embeddings).astype("float32")
    index = faiss.IndexFlatL2(384)
    index.add(embeddings_array)
    return index

def retrieve_context(query, top_k=3):
    if st.session_state.faiss_index is None: return []
    query_vector = np.array(embed_model.encode([query])).astype("float32")
    _, indices = st.session_state.faiss_index.search(query_vector, top_k)
    return [st.session_state.text_chunks[idx] for idx in indices[0] if idx != -1 and idx < len(st.session_state.text_chunks)]

def extract_text_from_file(uploaded_file, file_extension):
    extracted_text = ""
    if file_extension == "pdf":
        pdf = PdfReader(uploaded_file)
        extracted_text = "\n".join([p.extract_text() for p in pdf.pages if p.extract_text()])
    elif file_extension == "docx":
        doc = Document(uploaded_file)
        extracted_text = "\n\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
    elif file_extension == "pptx":
        prs = Presentation(uploaded_file)
        slide_texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
        extracted_text = "\n".join(slide_texts)
    return extracted_text

# ==========================================
# 1. SIDEBAR PANEL: CONVERSATION ARCHIVE
# ==========================================
with st.sidebar:
    st.title("📜 Conversation Logs")
    st.caption("Active Session History Panel")
    st.divider()
    
    st.subheader("📄 Saved Document Chats")
    if st.button("➕ New Document Chat", key="btn_new_doc_chat"):
        if st.session_state.chat_history:
            session_title = f"Doc Session #{len(st.session_state.saved_doc_sessions) + 1}: {st.session_state.chat_history[0]['content'][:20]}..."
            st.session_state.saved_doc_sessions[session_title] = st.session_state.chat_history
        st.session_state.chat_history = []
        st.rerun()
        
    for title, history in st.session_state.saved_doc_sessions.items():
        if st.button(f"💬 {title}", key=f"load_doc_{title}"):
            st.session_state.chat_history = history
            st.rerun()
            
    st.divider()
    
    st.subheader("🖼️ Saved Vision Chats")
    if st.button("➕ New Vision Chat", key="btn_new_img_chat"):
        if st.session_state.image_chat_history:
            session_title = f"Vision Session #{len(st.session_state.saved_img_sessions) + 1}: {st.session_state.image_chat_history[0]['content'][:20]}..."
            st.session_state.saved_img_sessions[session_title] = st.session_state.image_chat_history
        st.session_state.image_chat_history = []
        st.rerun()
        
    for title, history in st.session_state.saved_img_sessions.items():
        if st.button(f"🖼️ {title}", key=f"load_img_{title}"):
            st.session_state.image_chat_history = history
            st.rerun()

# ==========================================
# 2. MAIN APP SPACE: TABBED WORKSPACE
# ==========================================
st.title("🤖 Universal Multimodal AI Workspace")
tab_doc, tab_img = st.tabs(["📄 Universal Document RAG", "🖼️ Vision Assistant Engine"])

# --- TAB 1: UNIVERSAL DOCUMENT RAG ---
with tab_doc:
    st.subheader("📁 Upload Document")
    uploaded_file = st.file_uploader("Upload Document (PDF, DOCX, PPTX)", type=["pdf", "docx", "pptx"], key="doc_file_uploader")
    
    file_id = f"{uploaded_file.name}_{uploaded_file.size}" if uploaded_file else None
    
    if uploaded_file:
        if st.session_state.get("current_file_id") != file_id:
            with st.status("Processing....", expanded=True) as status:
                file_ext = uploaded_file.name.split(".")[-1].lower()
                st.session_state.current_file_ext = file_ext
                st.session_state.current_file_name = uploaded_file.name
                st.session_state.file_raw_bytes = uploaded_file.getvalue()
                
                extracted_text = extract_text_from_file(uploaded_file, file_ext)
                if extracted_text.strip():
                    st.session_state.doc_raw_text = extracted_text
                    chunks = chunk_text(extracted_text)
                    st.session_state.text_chunks = chunks
                    st.session_state.faiss_index = build_vector_db(chunks)
                    st.session_state.current_file_id = file_id
                    st.session_state.chat_history = []
                    if "doc_summary" in st.session_state: del st.session_state.doc_summary
                    status.update(label=f"Successfully {file_ext.upper()} completed...", state="complete")
                else:
                    status.update(label="Failed to extract layout data.", state="error")
                    st.warning("Could not read any structural text from this document layout.")
                    
    if st.session_state.doc_raw_text:
        col_btn1, col_btn2, _ = st.columns([1.2, 1, 2])
        
        # --- SELECTIVE CONDITIONAL PREVIEW LOGIC ---
        is_pdf = st.session_state.current_file_ext == "pdf"
        
        with col_btn1:
            if is_pdf:
                show_preview = st.toggle("🔍PDF Document Viewer", value=False)
            else:
                show_preview = False
                # Custom Metadata Asset Block for office types with a built-in launcher button
                st.download_button(
                    label=f"📥 Download to open in MS { 'Word' if st.session_state.current_file_ext == 'docx' else 'PowerPoint' }",
                    data=st.session_state.file_raw_bytes,
                    file_name=st.session_state.current_file_name,
                    mime="application/octet-stream"
                )
        with col_btn2:
            trigger_summary = st.button("📋 Generate Executive Summary", type="primary")
            
        # PDF Preview Window block
        if is_pdf and show_preview and st.session_state.file_raw_bytes:
            with st.container(border=True):
                pdf_viewer(st.session_state.file_raw_bytes)
                
        if trigger_summary:
            with st.spinner("Generating overview summary..."):
                headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
                prompt = f"Summarize the following document concisely into Main Topics, Key Points,short description and Takeaways:\n\n{st.session_state.doc_raw_text[:6000]}"
                response = requests.post("https://openrouter.ai/api/v1/chat/completions", headers=headers, json={"model": "openrouter/free", "messages": [{"role": "user", "content": prompt}]})
                result = response.json()
                if "choices" in result:
                    st.session_state.doc_summary = result["choices"][0]["message"]["content"]

        if "doc_summary" in st.session_state:
            with st.container(border=True):
                st.markdown("### 📋 Executive Summary")
                st.markdown(st.session_state.doc_summary)
                if st.button("❌ Close Summary View"):
                    del st.session_state.doc_summary
                    st.rerun()

        st.divider()
        st.subheader("💬 Interactive Context RAG Conversation")
        
        for message in st.session_state.chat_history:
            with st.chat_message(message["role"]):
                st.write(message["content"])
        
        if user_question := st.chat_input("Ask a question about the active document...", key="doc_chat_input"):
            with st.chat_message("user"):
                st.write(user_question)
            
            with st.spinner("Analyzing vector index..."):
                matched_chunks = retrieve_context(user_question, top_k=3)
                context_str = "\n---\n".join(matched_chunks)
                
                system_guideline = (
                    "You are a helpful document assistant. Answer the user's question using ONLY the provided text context. "
                    "If the answer cannot be derived from the text, state clearly that the document does not contain the answer. "
                    f"\n\nGROUNDED DOCUMENT CONTEXT FOR THIS TURN:\n{context_str}"
                )
                
                api_messages = [{"role": "system", "content": system_guideline}]
                for msg in st.session_state.chat_history:
                    api_messages.append({"role": msg["role"], "content": msg["content"]})
                api_messages.append({"role": "user", "content": user_question})
                
                headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
                response = requests.post("https://openrouter.ai/api/v1/chat/completions", headers=headers, json={"model": "openrouter/free", "messages": api_messages})
                result = response.json()
                
                if "choices" in result:
                    answer = result["choices"][0]["message"]["content"]
                    with st.chat_message("assistant"):
                        st.write(answer)
                        
                    st.session_state.chat_history.append({"role": "user", "content": user_question})
                    st.session_state.chat_history.append({"role": "assistant", "content": answer})
                    st.rerun()

# --- TAB 2: VISION ASSISTANT ---
with tab_img:
    st.subheader("📁 Upload Workspace Image")
    image_file = st.file_uploader("Drop Image File Here", type=["png", "jpg", "jpeg"], key="img_file_uploader")
    
    if image_file:
        image = Image.open(image_file)
        img_bytes = image_file.getvalue()
        encoded_img = base64.b64encode(img_bytes).decode("utf-8")
        if st.session_state.active_image_b64 != encoded_img:
            st.session_state.active_image_b64 = encoded_img
            st.session_state.image_chat_history = []
            
        # --- NATIVE DEDICATED IMAGE PREVIEW FRAME ---
        show_img_preview = st.toggle("🔍 Image Preview", value=True)
        if show_img_preview:
            st.image(image, caption="Active Image File Workspace", width=350)
            
        st.divider()
        st.subheader("💬 Chat about Image")
        
        for msg in st.session_state.image_chat_history:
            with st.chat_message(msg["role"]):
                st.write(msg["content"])
                
        if image_question := st.chat_input("Ask something about this image...", key="img_chat_input"):
            with st.chat_message("user"):
                st.write(image_question)
                
            with st.spinner("Analyzing image assets..."):
                headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
                
                api_messages = [
                    {
                        "role": "system",
                        "content": "You are an advanced computer vision assistant. Have a continuous discussion with the user regarding the provided image input."
                    }
                ]
                for msg in st.session_state.image_chat_history:
                    api_messages.append({"role": msg["role"], "content": msg["content"]})
                    
                api_messages.append({
                    "role": "user",
                    "content": [
                        {"type": "text", "text": image_question},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{st.session_state.active_image_b64}"}}
                    ]
                })
                
                response = requests.post(
                    "https://openrouter.ai/api/v1/chat/completions",
                    headers=headers,
                    json={"model": "openrouter/free", "messages": api_messages}
                )
                result = response.json()
                
                if "choices" in result:
                    answer = result["choices"][0]["message"]["content"]
                    with st.chat_message("assistant"):
                        st.write(answer)
                        
                    st.session_state.image_chat_history.append({"role": "user", "content": image_question})
                    st.session_state.image_chat_history.append({"role": "assistant", "content": answer})
                    st.rerun()
                else:
                    st.error("Vision API pipeline failure.")
    else:
        st.info("Please upload an image matrix file to enable computer vision operations.")
